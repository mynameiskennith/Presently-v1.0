# import streamlit as st
from fastapi.responses import JSONResponse
import pptx
from pptx.util import Inches, Pt
import os
# from dotenv import load_dotenv
from groq import Groq
from io import BytesIO
import tempfile
import zipfile
import xml.etree.ElementTree as ET
# import time
from PIL import Image as Imagee
# import sounddevice as sd
# import numpy as np
# import scipy.io.wavfile as wavfile
# import speech_recognition as sr
from spire.presentation import Presentation
from spire.presentation.common import *
# import wave
# import struct
import json
import traceback
from io import BytesIO
from pptx import Presentation



# Get API key from environment variable
api_key = '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'
if not api_key:
    raise ValueError("GROQ_API_KEY not found in environment variables")

client = Groq(api_key=api_key)

# Custom Formatting options
TITLE_FONT_SIZE = Pt(30)
SLIDE_FONT_SIZE = Pt(16)

def generate_slide_titles(topic,noOfSlides,audienceType,slideContent):
    prompt = f"""Generate exactly {noOfSlides} concise slide titles for a presentation on the topic: {topic}
    Rules:
    1. Provide only the titles, one per line
    2. Do not include any numbering or bullet points
    3. Each title should be brief and relevant to the topic
    4. Do not include any additional text in response  or explanations
    6. Directly give the titles and do not add any additional message above it
    7. The type of presentation is {audienceType} 
    8. Give more importance to : {slideContent} when selecting topics """
    
    response = client.chat.completions.create(
        model="llama3-8b-8192",
        messages=[
            {
                "role": "user",
                "content": prompt
            }
        ],
        temperature=0.7,
        max_tokens=1024,
        top_p=1,
        stream=False,
    )

    # Extract the content from the response
    response_text = response.choices[0].message.content

    # Split the response into titles and filter out empty lines
    return [title.strip() for title in response_text.split("\n") if title.strip()] 



def generate_slide_content(slide_title,audienceType):
    prompt = f"""Generate exactly 7 bullet points for the slide titled: "{slide_title}"
    Rules:
    1. Each point must be a very short but crisp sentence
    2. Do not exceed 15 words per point
    3. Provide only the points, one per line
    4. Do not include any numbering or bullet point symbols
    5. Do not include any additional text from response or 
    6. Each point should be self explanatory
    7. Directly provide the points for the slide title and do not include any additional message before the points
    8. Do not include the slide title in the points
    9. The type of presentation is {audienceType} , select the tone of points accordingly"""
    
    response = client.chat.completions.create(
        model="llama3-8b-8192",
        messages=[
            {
                "role": "user",
                "content": prompt
            }
        ],
        temperature=0.7,
        max_tokens=1024,
        top_p=1,
        stream=False,
    )

    # Extract the content from the response
    response_text = response.choices[0].message.content

    # Split the response into points and filter out empty lines
    points = [point.strip() for point in response_text.split("\n") if point.strip()][1:7]  # Ensure we get exactly 6 points
    
    # Join the points with newlines to create the slide content
    return "\n\n".join(points)



def create_presentation(request_data):
    slide_titles = generate_slide_titles(request_data["topic"],request_data["numberOfSlides"],request_data["audienceType"],request_data["slideContent"])
    slide_contents = []
    for title in slide_titles:
        print(f'Generating slide {title}')
        slide_content = generate_slide_content(title,request_data["audienceType"])
        slide_contents.append(slide_content)

    prs = pptx.Presentation()
    slide_layout = prs.slide_layouts[1]

    # Add the title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = request_data["topic"]
    print("Added the title slide")

    # Add the "Contents" slide (second slide)
    contents_slide = prs.slides.add_slide(slide_layout)
    contents_slide.shapes.title.text = "Contents"
    print("Added the contents title")

    # Add each slide title as bullet points on the Contents slide
    content_text = "\n\n".join(slide_titles)  # Join all slide titles into one text with line breaks
    contents_slide.shapes.placeholders[1].text = content_text
    print("Added the contents")


    # Customize font size for the "Contents" slide
    contents_slide.shapes.title.text_frame.paragraphs[0].font.size = TITLE_FONT_SIZE
    for paragraph in contents_slide.shapes.placeholders[1].text_frame.paragraphs:
        paragraph.font.size = SLIDE_FONT_SIZE
    print("Customized the font size")

    # Add the rest of the slides with their respective content
    for slide_title, slide_content in zip(slide_titles, slide_contents):
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = slide_title
        slide.shapes.placeholders[1].text = slide_content
        
        # Customize font size for each slide
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    paragraph.font.size = SLIDE_FONT_SIZE

        slide.shapes.title.text_frame.paragraphs[0].font.size = TITLE_FONT_SIZE
        print("Created A Slide")
                    
                    
    thank_you_slide = prs.slides.add_slide(prs.slide_layouts[0])
    thank_you_slide.shapes.title.text = "Thank You"
    print("Cretaed the Thank You")

    # Save to an in-memory byte stream
    ppt_stream = BytesIO()
    prs.save(ppt_stream)
    ppt_stream.seek(0)  # Reset stream position for reading
    return ppt_stream

    # Save the presentation
    # os.makedirs('generated_ppt', exist_ok=True)
    # ppt_path = os.path.join('generated_ppt', f'{request_data["topic"]}_presentation.pptx')
    # print(f"////////////////////\n////////////////\n{ppt_path}")
    # prs.save(ppt_path)
    # return ppt_path



# def rate_ppt(ppt_file_contents):
#     with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp_file:
#         tmp_file.write(ppt_file_contents)
#         tmp_file_path = tmp_file.name

#     try:
#         with zipfile.ZipFile(tmp_file_path, 'r') as zip_ref:
#             # Get total number of slides
#             presentation_xml = zip_ref.read('ppt/presentation.xml')
#             root = ET.fromstring(presentation_xml)
#             total_slides = len(root.findall('.//{http://schemas.openxmlformats.org/presentationml/2006/main}sldId'))

#             # Initialize ratings
#             slide_number_rating = 10
#             bullet_point_rating = 10

#             # Analyze each slide
#             for i in range(1, total_slides + 1):
#                 try:
#                     slide_xml = zip_ref.read(f'ppt/slides/slide{i}.xml')
#                     slide_root = ET.fromstring(slide_xml)

#                     # Check for slide number
#                     slide_number = slide_root.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}ph[@type="sldNum"]')
#                     if slide_number is None:
#                         slide_number_rating -= 0.5

#                     # Count bullet points
#                     bullet_points = slide_root.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}buChar')
#                     if len(bullet_points) > 7:
#                         bullet_point_rating -= 1

#                 except KeyError:
#                     # Slide doesn't exist, reduce rating
#                     slide_number_rating -= 1

#         # Ensure ratings don't go below 0
#         slide_number_rating = max(0, slide_number_rating)
#         bullet_point_rating = max(0, bullet_point_rating)

#         # Calculate overall rating
#         overall_rating = (slide_number_rating + bullet_point_rating) / 2

#         return {
#             "overall_rating": overall_rating,
#             "slide_number_rating": slide_number_rating,
#             "bullet_point_rating": bullet_point_rating,
#             "total_slides": total_slides
#         }
    
#     except Exception as error:
#     # handle the exception
#         return {"Error" : error}

#     finally:
#         # Clean up the temporary file
#         if os.path.exists(tmp_file_path):
#             os.unlink(tmp_file_path)


def extract_text_from_ppt(ppt_file_contents):
    """Extract all text from the PPT for processing."""
    presentation = Presentation(BytesIO(ppt_file_contents))
    all_text = []
    for slide in presentation.slides:
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                slide_text.append(shape.text)
        all_text.append("\n".join(slide_text))
    return "\n\n".join(all_text)

def rate_ppt(ppt_file_contents):
    """Evaluate a PPT file using the Groq model."""
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp_file:
        tmp_file.write(ppt_file_contents)
        tmp_file_path = tmp_file.name

    try:
        # Extract text from the PPT for the prompt
        ppt_text = extract_text_from_ppt(ppt_file_contents)

        # Prepare the prompt for the Groq model
        prompt = f"""
        You are an AI assistant tasked with evaluating PowerPoint presentations.
        
        PPT content: {ppt_text}
        
        Analyze the provided presentation based on these criteria:
        1. Number of points per slide (score out of 10).
        2. Number of images per slide (score out of 10).
        3. Readability of text content - SMOG Readability (score out of 10).
        4. Consistency of slide formatting (score out of 10).
        5. Overall presentation content quality (score out of 10).
        6. Number of Slides (score out of 10).
        7. Overall score. (score out of 10)

        Return a JSON object in this format:
        {{
            "noOfPoints": score,
            "noOfImages": score,
            "Readability": score,
            "Consistency": score,
            "Quality": score,
            "noOfSlides": score,
            "overAllScore": score
        }}

        Only return the JSON object with no additional text.
        """

        # Pass the prompt to the Groq model
        response = client.chat.completions.create(
            model="llama3-8b-8192",
            messages=[
                {
                    "role": "user",
                    "content": prompt
                }
            ],
            temperature=0.7,
            max_tokens=1024,
            top_p=1,
            stream=False,
        )

        # Extract and parse the response
        evaluation_results = response.choices[0].message.content.strip()
        try:
            evaluation_results = json.loads(evaluation_results)
        except json.JSONDecodeError:
            raise ValueError("Invalid JSON response from Groq API")

        return evaluation_results

    except Exception as error:
        print(traceback.format_exc())  # Log detailed error
        return {"Error": str(error)}

    finally:
        if os.path.exists(tmp_file_path):
            os.unlink(tmp_file_path)



# Function to convert PPT to images
def convert_ppt_to_images(ppt_file):
    # Save the uploaded file to a temporary file
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp_ppt_file:
        tmp_ppt_file.write(ppt_file.getvalue())  # Write the content of the uploaded file
        tmp_ppt_file_path = tmp_ppt_file.name    # Get the file path

    try:
        # Load the presentation from the saved file
        presentation = Presentation()
        presentation.LoadFromFile(tmp_ppt_file_path)

        images = []
        with tempfile.TemporaryDirectory() as tmpdirname:
            for i, slide in enumerate(presentation.Slides):
                fileName = f"slide_{i}.png"
                image = slide.SaveAsImageByWH(800, 450)
                file_path = os.path.join(tmpdirname, fileName)
                image.Save(file_path)
                image.Dispose()

                # Open the image, convert it to RGB (to ensure compatibility), and store it in memory
                with Imagee.open(file_path) as img:
                    images.append(img.copy().convert('RGB'))

        presentation.Dispose()
        return {"images":images}

    finally:
        # Clean up the temporary file
        if os.path.exists(tmp_ppt_file_path):
            os.remove(tmp_ppt_file_path)



# Function to transcribe audio using SpeechRecognition
# def transcribe_audio(audio_file):
    # recognizer = sr.Recognizer()
    # with sr.AudioFile(audio_file) as source:
    #     audio_data = recognizer.record(source)
    #     try:
    #         text = recognizer.recognize_google(audio_data)
    #         return text
    #     except sr.UnknownValueError:
    #         return "Audio is unclear. Could not transcribe."
    #     except sr.RequestError as e:
    #         return f"Could not request results; {e}"